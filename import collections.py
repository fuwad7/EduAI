import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Initialize Presentation
prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9 Widescreen standard
prs.slide_height = Inches(7.5)

# --- THEME DESIGN SYSTEM ---
COLOR_BG = RGBColor(7, 10, 14)       # Ultra Dark Charcoal Navy
COLOR_TEAL = RGBColor(43, 196, 182)  # Glowing Teal
COLOR_GOLD = RGBColor(217, 119, 6)   # Deep Gold Accent
COLOR_WHITE = RGBColor(255, 255, 255)# Crisp Text White
COLOR_GRAY = RGBColor(148, 163, 184) # Muted Blueprint Gray

FONT_TITLE = 'Georgia'   # Elegant serif alternative for Cinzel
FONT_BODY = 'Arial'      # Clean, high-readability sans-serif

def set_background(slide):
    """Fills slide background with a premium deep dark tint."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG

def add_header(slide, title_text, subtitle_text=""):
    """Creates a consistent, visually catching top header banner."""
    # Top accent line
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.08)) # 1 = Rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_TEAL
    shape.line.fill.background()

    # Title box
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    p = tf.paragraphs[0]
    p.text = title_text.upper()
    p.font.name = FONT_TITLE
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    
    if subtitle_text:
        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        p2.font.name = FONT_BODY
        p2.font.size = Pt(14)
        p2.font.color.rgb = COLOR_GOLD
        p2.space_before = Pt(4)

def add_footer(slide, slide_num):
    """Appends institutional baseline identifier tags."""
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(6.9), Inches(11.7), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"DIU EduAI Platform  |  Academic Operations Matrix"
    p.font.name = FONT_BODY
    p.font.size = Pt(10)
    p.font.color.rgb = COLOR_GRAY
    
    # Slide Number Alignment Right
    p_num = tf.add_paragraph()
    p_num.text = str(slide_num).zfill(2)
    p_num.alignment = PP_ALIGN.RIGHT
    p_num.font.name = FONT_BODY
    p_num.font.size = Pt(10)
    p_num.font.bold = True
    p_num.font.color.rgb = COLOR_TEAL

def create_card(slide, left, top, width, height, title, body, is_gold=False):
    """Creates streamlined visual placeholder cards for text containment."""
    card = slide.shapes.add_shape(1, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(16, 22, 30) # Soft card elevation color
    # Subtle accent border line
    card.line.color.rgb = COLOR_GOLD if is_gold else COLOR_TEAL
    card.line.width = Pt(1.5)
    
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.25)
    
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = FONT_TITLE
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_GOLD if is_gold else COLOR_TEAL
    p.space_after = Pt(8)
    
    p2 = tf.add_paragraph()
    p2.text = body
    p2.font.name = FONT_BODY
    p2.font.size = Pt(13)
    p2.font.color.rgb = COLOR_GRAY

blank_layout = prs.slide_layouts[6]

# ==============================================================================
# SLIDE 1: TITLE SLIDE
# ==============================================================================
slide1 = prs.slides.add_slide(blank_layout)
set_background(slide1)

# Centralized framework banner accent
mid_bar = slide1.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.1))
mid_bar.fill.solid()
mid_bar.fill.fore_color.rgb = COLOR_TEAL
mid_bar.line.fill.background()

t_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(3.5))
tf1 = t_box.text_frame
tf1.word_wrap = True

p1 = tf1.paragraphs[0]
p1.alignment = PP_ALIGN.CENTER
p1.text = "DIU EduAI PLATFORM"
p1.font.name = FONT_TITLE
p1.font.size = Pt(56)
p1.font.bold = True
p1.font.color.rgb = COLOR_WHITE

p2 = tf1.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
p2.text = "Advanced AI Teaching, Assessment & Curriculum Portal"
p2.font.name = FONT_TITLE
p2.font.size = Pt(22)
p2.font.color.rgb = COLOR_TEAL
p2.space_before = Pt(12)

p3 = tf1.add_paragraph()
p3.alignment = PP_ALIGN.CENTER
p3.text = "DAFFODIL INTERNATIONAL UNIVERSITY  •  PROJECT DEFENSE"
p3.font.name = FONT_BODY
p3.font.size = Pt(12)
p3.font.color.rgb = COLOR_GRAY
p3.space_before = Pt(60)

add_footer(slide1, 1)

# ==============================================================================
# SLIDE 2: EXECUTIVE SUMMARY
# ==============================================================================
slide2 = prs.slides.add_slide(blank_layout)
set_background(slide2)
add_header(slide2, "1. Executive Summary", "The Academic Co-Pilot Overview")

# Narrative Statement Left Box
s2_box = slide2.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(7.0), Inches(2.5))
tf2 = s2_box.text_frame
tf2.word_wrap = True
p_s2 = tf2.paragraphs[0]
p_s2.text = "A unified, AI-driven gateway engineered to automate exhausting institutional workflows, maintain strict data safety, and unify loose academic tools."
p_s2.font.name = FONT_BODY
p_s2.font.size = Pt(24)
p_s2.font.color.rgb = COLOR_WHITE

# Sub-context fields
p_sub1 = tf2.add_paragraph()
p_sub1.text = "Target Ecosystem: Faculty, Chairs, & External Program Auditors.\nCompliance Focus: Native alignment with BAET / ABET accreditation records."
p_sub1.font.size = Pt(14)
p_sub1.font.color.rgb = COLOR_GRAY
p_sub1.space_before = Pt(30)

# Giant Metric Badge Right
badge = slide2.shapes.add_shape(1, Inches(8.5), Inches(2.2), Inches(4.0), Inches(3.8))
badge.fill.solid()
badge.fill.fore_color.rgb = RGBColor(11, 28, 26)
badge.line.color.rgb = COLOR_TEAL
btf = badge.text_frame
btf.word_wrap = True
btf.margin_top = Inches(0.6)

bp1 = btf.paragraphs[0]
bp1.alignment = PP_ALIGN.CENTER
bp1.text = "10+"
bp1.font.name = FONT_TITLE
bp1.font.size = Pt(80)
bp1.font.bold = True
bp1.font.color.rgb = COLOR_TEAL

bp2 = btf.add_paragraph()
bp2.alignment = PP_ALIGN.CENTER
bp2.text = "HOURS RECLAIMED / WEEK\n\nAutomating routine processing loops without compromising core teaching rigor."
bp2.font.name = FONT_BODY
bp2.font.size = Pt(13)
bp2.font.color.rgb = COLOR_WHITE

add_footer(slide2, 2)

# ==============================================================================
# SLIDE 3: PROBLEM STATEMENT (Simplified Matrix Layout)
# ==============================================================================
slide3 = prs.slides.add_slide(blank_layout)
set_background(slide3)
add_header(slide3, "2. Institutional Challenges", "Operational Roadblocks & Consequences")

# Grid matrix mapping out 4 pain points neatly
problems = [
    ("Administrative Overload", "Hours lost on grading scripts & rosters.", "Reduced student mentorship focus."),
    ("OBE Matrix Complexity", "Error-prone manual outcome mapping tracking.", "Accreditation audit vulnerability."),
    ("Fragmented Tools", "Educators constantly shifting across disjointed apps.", "Severe organizational data silos."),
    ("Capstone Supervision", "No centralized milestone visibility layout structures.", "Unpredictable defense delays.")
]

for i, (title, impact, consequence) in enumerate(problems):
    x = Inches(0.8 + (i % 2) * 6.0)
    y = Inches(2.0 + (i // 2) * 2.3)
    create_card(slide3, x, y, Inches(5.6), Inches(2.0), f"✕  {title}", f"System Impact: {impact}\nConsequence: {consequence}")

add_footer(slide3, 3)

# ==============================================================================
# SLIDE 4: VALUE PROPOSITION
# ==============================================================================
slide4 = prs.slides.add_slide(blank_layout)
set_background(slide4)
add_header(slide4, "3. Value Proposition", "Consolidating Academic Operations")

props = [
    ("Unified Command", "One single interface portal tracking analytics and metrics safely.", False),
    ("Intelligent AI", "Instant test optimization, dynamic evaluation matrix generation.", False),
    ("OBE Native Design", "Outcome trails built straight into tables for clean audit records.", False),
    ("Row-Level Safety", "Postgres policy rules preventing cross-user profile exposure.", True),
    ("Offline Buffer", "LocalStorage syncing maintaining lightning fast page caching loops.", True)
]

for i, (title, desc, gold) in enumerate(props):
    x = Inches(0.8 + (i % 3) * 3.9)
    y = Inches(2.0 + (i // 3) * 2.3)
    create_card(slide4, x, y, Inches(3.6), Inches(2.0), title, desc, is_gold=gold)

# Closing accent block
extra = slide4.shapes.add_shape(1, Inches(8.6), Inches(4.3), Inches(3.6), Inches(2.0))
extra.fill.solid()
extra.fill.fore_color.rgb = RGBColor(24, 18, 10)
extra.line.color.rgb = COLOR_GOLD
etf = extra.text_frame
etf.word_wrap = True
ep = etf.paragraphs[0]
ep.alignment = PP_ALIGN.CENTER
ep.text = "\nSHIFTING FOCUS TO ACTIVE STRATEGIC MENTORSHIP"
ep.font.name = FONT_TITLE
ep.font.size = Pt(14)
ep.font.bold = True
ep.font.color.rgb = COLOR_GOLD

add_footer(slide4, 4)

# ==============================================================================
# SLIDE 5: TECHNICAL INFRASTRUCTURE
# ==============================================================================
slide5 = prs.slides.add_slide(blank_layout)
set_background(slide5)
add_header(slide5, "4. Technical Infrastructure", "The Platform Technology Stack")

create_card(slide5, Inches(0.8), Inches(2.2), Inches(5.6), Inches(4.2), 
            "Client & Display Layers", 
            "• React.js Framework: Pure functional hooks configuration for responsive UI updates.\n\n• Low-Fatigue Interface: Custom academic dark motif minimizing visual strain during grading.\n\n• LocalStorage Mirror: Instant edge query caching ensuring fluid navigation loops offline.")

create_card(slide5, Inches(6.9), Inches(2.2), Inches(5.6), Inches(4.2), 
            "Data & Shield Environments", 
            "• Supabase (PostgreSQL): Secure relational storage managing core schemas reliably.\n\n• Isolated Storage Buckets: Private sandbox vaults processing heavy source scripts and PDF data.\n\n• Row-Level Hardening (RLS): Multi-tenant isolation verified directly inside the database engine.", True)

add_footer(slide5, 5)

# ==============================================================================
# SLIDE 6: AI ENGINE PLATFORM
# ==============================================================================
slide6 = prs.slides.add_slide(blank_layout)
set_background(slide6)
add_header(slide6, "5. AI Engine Infrastructure", "Fault-Tolerant Dynamic Routing System")

create_card(slide6, Inches(0.8), Inches(2.2), Inches(3.6), Inches(4.2), 
            "Tier 01: Groq API", 
            "Primary pipeline managing speed-critical processes using Llama-3.3-70b.\n\nOptimized for rapid generation of quiz iterations, custom notification briefs, and correspondence layouts.")

create_card(slide6, Inches(4.8), Inches(2.2), Inches(3.6), Inches(4.2), 
            "Tier 02: Google Gemini", 
            "Secondary core processing reasoning tasks using Gemini-2.5-flash.\n\nLeveraged for deep semantic context analytics, syllabus diagnostic mapping, and grading rubrics matrices.")

create_card(slide6, Inches(8.8), Inches(2.2), Inches(3.7), Inches(4.2), 
            "Tier 03: OpenRouter Backup", 
            "Safeguard edge route running GPT-4o-mini to guarantee 100% platform uptime.\n\nInvariants are locked within environmental prompts enforcing an immutable academic persona.", True)

add_footer(slide6, 6)

# ==============================================================================
# SLIDE 7: CORE MODULES OVERVIEW (Combined & Simplified)
# ==============================================================================
slide7 = prs.slides.add_slide(blank_layout)
set_background(slide7)
add_header(slide7, "6. Core Functional Portals", "Faculty Command Center Capabilities")

modules = [
    ("Faculty Hub", "Sub-millisecond global search crossing tracking metrics, portfolios, and courses dynamically."),
    ("Assessment Engine", "Automated MCQ creation mapped to learning outcomes, with integrated plagiarism pattern checkers."),
    ("OBE Matrix Map", "Accreditation rule compliance tools linking assignments to strict BAET / ABET review nodes."),
    ("Syllabus Engine", "Transforms flat reading assets into direct structural guides, version histories, and interactive FAQs.")
]

for i, (title, desc) in enumerate(modules):
    x = Inches(0.8 + (i % 2) * 6.0)
    y = Inches(2.2 + (i // 2) * 2.2)
    create_card(slide7, x, y, Inches(5.6), Inches(1.9), title, desc)

add_footer(slide7, 7)

# ==============================================================================
# SLIDE 8: CAPSTONE PROJECT TRACKING
# ==============================================================================
slide8 = prs.slides.add_slide(blank_layout)
set_background(slide8)
add_header(slide8, "7. Capstone Supervision", "Standardized Project Pipeline Tracking")

# Visual Linear Pipeline Layout Nodes
steps = [
    ("01 / Proposal", "AI evaluates core data logic for resource feasibility."),
    ("02 / Lit Review", "Heuristic crosscheck parses early abstract anomalies."),
    ("03 / Build Loop", "Telemetry dashboards record continuous progression."),
    ("04 / Log Diary", "Advisors lock precise student action items instantly."),
    ("05 / Final Jury", "Automated system loops enforce defense accountability.")
]

for i, (step_title, step_desc) in enumerate(steps):
    x = Inches(0.8 + i * 2.4)
    y = Inches(2.8)
    # Highlight final phases in gold
    gold_status = i >= 3
    create_card(slide8, x, y, Inches(2.2), Inches(3.4), step_title, step_desc, is_gold=gold_status)

add_footer(slide8, 8)

# ==============================================================================
# SLIDE 9: STRATEGIC QUANTIFIED IMPACT
# ==============================================================================
slide9 = prs.slides.add_slide(blank_layout)
set_background(slide9)
add_header(slide9, "8. Quantified Impact Target", "Strategic Metrics and Institutional Stakeholder Yields")

impacts = [
    ("Faculty Members", "Automates routine tasks, reclaiming 10+ hours weekly for advanced academic research."),
    ("Student Cohorts", "Delivers prompt evaluation turnarounds paired with clear, transparent rubric breakdowns."),
    ("Administration", "Maintains living outcome evidence streams, dramatically cutting audit preparation timelines."),
    ("Department Chairs", "Provides high-level real-time visibility over final-year project completion velocity.")
]

for i, (stakeholder, detail) in enumerate(impacts):
    x = Inches(0.8 + (i % 2) * 6.0)
    y = Inches(2.2 + (i // 2) * 2.2)
    create_card(slide9, x, y, Inches(5.6), Inches(1.9), stakeholder, detail, is_gold=(i >= 2))

add_footer(slide9, 9)

# ==============================================================================
# SLIDE 10: STRATEGIC ROADMAP
# ==============================================================================
slide10 = prs.slides.add_slide(blank_layout)
set_background(slide10)
add_header(slide10, "9. Strategic Roadmap", "Future Expansion Milestones")

roadmaps = [
    ("LMS Native Synced Engines", "Deep API integrations supporting direct gradebook connectivity across Canvas and Moodle frameworks."),
    ("Voice Advisory Summarizers", "Automatic acoustic transcription processing supervisor records into actionable follow-up items."),
    ("Predictive Predictive Analytics", "Machine learning diagnostics tracking metrics to flag vulnerable students prior to exam phases."),
    ("AI Peer Moderation", "Decentralized anonymous evaluation nodes running built-in anomaly tracking filtering engines.")
]

for i, (title, desc) in enumerate(roadmaps):
    x = Inches(0.8 + i * 2.95)
    y = Inches(2.4)
    create_card(slide10, x, y, Inches(2.75), Inches(4.0), title, desc)

add_footer(slide10, 10)

# ==============================================================================
# SLIDE 11: CONCLUSION & QA
# ==============================================================================
slide11 = prs.slides.add_slide(blank_layout)
set_background(slide11)

t_box11 = slide11.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(3.5))
tf11 = t_box11.text_frame
tf11.word_wrap = True

p11_1 = tf11.paragraphs[0]
p11_1.alignment = PP_ALIGN.CENTER
p11_1.text = "THANK YOU"
p11_1.font.name = FONT_TITLE
p11_1.font.size = Pt(64)
p11_1.font.bold = True
p11_1.font.color.rgb = COLOR_WHITE

p11_2 = tf11.add_paragraph()
p11_2.alignment = PP_ALIGN.CENTER
p11_2.text = "DIU EduAI: Empowering Academia through Intelligent Automation"
p11_2.font.name = FONT_TITLE
p11_2.font.size = Pt(20)
p11_2.font.color.rgb = COLOR_TEAL
p11_2.space_before = Pt(12)

# Break line accent
line = slide11.shapes.add_shape(1, Inches(5.666), Inches(4.5), Inches(2.0), Inches(0.03))
line.fill.solid()
line.fill.fore_color.rgb = COLOR_GOLD
line.line.fill.background()

p11_3 = tf11.add_paragraph()
p11_3.alignment = PP_ALIGN.CENTER
p11_3.text = "QUESTIONS & ANSWERS"
p11_3.font.name = FONT_BODY
p11_3.font.size = Pt(16)
p11_3.font.bold = True
p11_3.font.color.rgb = COLOR_WHITE
p11_3.space_before = Pt(40)

add_footer(slide11, 11)

# Save presentation output
prs.save("DIU_EduAI_Defense.pptx")
print("Successfully generated 'DIU_EduAI_Defense.pptx'!")